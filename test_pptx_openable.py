#!/usr/bin/env python3
"""
Test if generated PPTX files can be opened and read by python-pptx.
This verifies that the files are truly openable and not corrupted.
"""

from pptx import Presentation
from pathlib import Path
import sys

def test_file(filepath):
    """Test if a PPTX file can be opened and read."""
    print(f"\n{'='*80}")
    print(f"Testing: {filepath}")
    print(f"{'='*80}")
    
    if not Path(filepath).exists():
        print(f"❌ File not found: {filepath}")
        return False
    
    try:
        # Try to open the file
        prs = Presentation(filepath)
        print(f"✅ File opened successfully")
        
        # Check properties
        print(f"\nProperties:")
        print(f"  - Slides: {len(prs.slides)}")
        print(f"  - Slide width: {prs.slide_width}")
        print(f"  - Slide height: {prs.slide_height}")
        
        # Check core properties
        try:
            props = prs.core_properties
            print(f"\nCore Properties:")
            print(f"  - Title: {props.title}")
            print(f"  - Author: {props.author}")
            print(f"  - Subject: {props.subject}")
        except Exception as e:
            print(f"  - Error reading properties: {e}")
        
        # Check slides
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1}:")
            print(f"  - Shapes: {len(slide.shapes)}")
            for j, shape in enumerate(slide.shapes):
                print(f"    - Shape {j+1}: {shape.shape_type} - {shape.name}")
                if hasattr(shape, "text"):
                    print(f"      Text: {shape.text[:50] if shape.text else '(empty)'}")
        
        print(f"\n✅ File is valid and readable")
        return True
        
    except Exception as e:
        print(f"❌ Error opening file: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    files = [
        "examples/output/simple_alignment_ppt_rs.pptx",
        "examples/output/enhanced_alignment_ppt_rs.pptx",
        "examples/output/reference_python_pptx.pptx",
    ]
    
    results = {}
    for filepath in files:
        results[filepath] = test_file(filepath)
    
    print(f"\n{'='*80}")
    print("Summary")
    print(f"{'='*80}")
    for filepath, success in results.items():
        status = "✅ PASS" if success else "❌ FAIL"
        print(f"{status}: {filepath}")
    
    all_pass = all(results.values())
    sys.exit(0 if all_pass else 1)
