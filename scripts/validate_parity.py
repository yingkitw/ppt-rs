#!/usr/bin/env python3
"""
Parity validation script - Compare ppt-rs output with python-pptx

This script generates PPTX files with both python-pptx and ppt-rs,
then compares them to ensure output parity.
"""

import zipfile
import xml.etree.ElementTree as ET
import sys
import os
from pathlib import Path
from pptx import Presentation

def normalize_xml(xml_string):
    """Normalize XML for comparison by removing whitespace differences."""
    root = ET.fromstring(xml_string)
    # Remove all text content (whitespace)
    for elem in root.iter():
        elem.text = None
        elem.tail = None
    return ET.tostring(root, encoding='unicode')

def compare_zip_structure(file1, file2):
    """Compare ZIP archive structure."""
    print("\n📁 Comparing ZIP Structure...")
    
    with zipfile.ZipFile(file1) as z1, zipfile.ZipFile(file2) as z2:
        files1 = sorted(z1.namelist())
        files2 = sorted(z2.namelist())
        
        if files1 != files2:
            print(f"❌ File list mismatch!")
            print(f"   python-pptx: {len(files1)} files")
            print(f"   ppt-rs: {len(files2)} files")
            
            only_in_1 = set(files1) - set(files2)
            only_in_2 = set(files2) - set(files1)
            
            if only_in_1:
                print(f"   Only in python-pptx: {only_in_1}")
            if only_in_2:
                print(f"   Only in ppt-rs: {only_in_2}")
            return False
        
        print(f"✅ ZIP structure matches ({len(files1)} files)")
        return True

def compare_xml_files(file1, file2, xml_paths):
    """Compare specific XML files."""
    print("\n📄 Comparing XML Content...")
    
    all_match = True
    with zipfile.ZipFile(file1) as z1, zipfile.ZipFile(file2) as z2:
        for xml_path in xml_paths:
            try:
                content1 = z1.read(xml_path)
                content2 = z2.read(xml_path)
                
                # Parse XML
                tree1 = ET.fromstring(content1)
                tree2 = ET.fromstring(content2)
                
                # Compare element counts
                count1 = len(list(tree1.iter()))
                count2 = len(list(tree2.iter()))
                
                if count1 == count2:
                    print(f"✅ {xml_path}: {count1} elements")
                else:
                    print(f"⚠️  {xml_path}: python-pptx={count1}, ppt-rs={count2}")
                    all_match = False
                    
            except KeyError:
                print(f"❌ {xml_path}: Missing in one file")
                all_match = False
    
    return all_match

def compare_file_sizes(file1, file2):
    """Compare file sizes."""
    print("\n📊 Comparing File Sizes...")
    
    size1 = os.path.getsize(file1)
    size2 = os.path.getsize(file2)
    
    diff = abs(size1 - size2)
    percent = (diff / max(size1, size2)) * 100
    
    print(f"python-pptx: {size1:,} bytes")
    print(f"ppt-rs:      {size2:,} bytes")
    print(f"Difference:  {diff:,} bytes ({percent:.1f}%)")
    
    # Allow up to 10% difference due to timestamps and metadata
    if percent <= 10:
        print("✅ File sizes are comparable")
        return True
    else:
        print("⚠️  File sizes differ significantly")
        return False

def validate_with_pptx(pptx_file):
    """Validate PPTX file with python-pptx."""
    print(f"\n🔍 Validating {Path(pptx_file).name} with python-pptx...")
    
    try:
        prs = Presentation(pptx_file)
        print(f"✅ File is valid")
        print(f"   Slides: {len(prs.slides)}")
        print(f"   Slide width: {prs.slide_width}")
        print(f"   Slide height: {prs.slide_height}")
        return True
    except Exception as e:
        print(f"❌ File is invalid: {e}")
        return False

def generate_python_pptx(output_file):
    """Generate PPTX with python-pptx."""
    print(f"\n🐍 Generating with python-pptx: {output_file}")
    
    try:
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(output_file)
        print(f"✅ Generated successfully")
        return True
    except Exception as e:
        print(f"❌ Generation failed: {e}")
        return False

def main():
    """Main validation function."""
    print("=" * 80)
    print("PARITY VALIDATION: ppt-rs vs python-pptx")
    print("=" * 80)
    
    # File paths
    python_file = "python_pptx_output.pptx"
    rust_file = "examples/output/01_simple.pptx"
    
    # Generate python-pptx file
    if not generate_python_pptx(python_file):
        sys.exit(1)
    
    # Verify rust file exists
    if not os.path.exists(rust_file):
        print(f"\n❌ ppt-rs file not found: {rust_file}")
        print("   Run: cargo run --example 01_create_simple_presentation")
        sys.exit(1)
    
    # Run comparisons
    results = []
    
    results.append(("ZIP Structure", compare_zip_structure(python_file, rust_file)))
    results.append(("XML Content", compare_xml_files(python_file, rust_file, [
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "docProps/core.xml"
    ])))
    results.append(("File Sizes", compare_file_sizes(python_file, rust_file)))
    results.append(("python-pptx validation", validate_with_pptx(python_file)))
    results.append(("ppt-rs validation", validate_with_pptx(rust_file)))
    
    # Summary
    print("\n" + "=" * 80)
    print("SUMMARY")
    print("=" * 80)
    
    passed = sum(1 for _, result in results if result)
    total = len(results)
    
    for test_name, result in results:
        status = "✅ PASS" if result else "❌ FAIL"
        print(f"{status}: {test_name}")
    
    print(f"\nTotal: {passed}/{total} tests passed")
    
    if passed == total:
        print("\n🎉 PARITY VERIFIED - ppt-rs output matches python-pptx!")
        return 0
    else:
        print("\n⚠️  Some tests failed - review differences above")
        return 1

if __name__ == "__main__":
    sys.exit(main())
